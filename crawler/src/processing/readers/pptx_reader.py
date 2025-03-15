"""
PPTX Document Reader

This module provides the PptxReader class for reading and extracting content from
Microsoft PowerPoint (.pptx) files.
"""

import os
from typing import Any, Optional

from .base import DocumentReader
from ..document_content import DocumentContent


class PptxReader(DocumentReader):
    """Reader for Microsoft PowerPoint (.pptx) files."""
    
    def read(self, file_path: str) -> DocumentContent:
        """Read PPTX file and extract content.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            DocumentContent object containing the extracted content
        """
        content = DocumentContent()
        try:
            from pptx import Presentation
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract basic metadata
            metadata = {
                'source': file_path,
                'format': 'pptx',
                'slide_count': len(prs.slides),
                'size_bytes': os.path.getsize(file_path)
            }
            
            # Try to extract more detailed properties if available
            try:
                core_props = prs.core_properties
                metadata.update({
                    'title': core_props.title or "",
                    'author': core_props.author or "",
                    'subject': core_props.subject or "",
                    'keywords': core_props.keywords or "",
                    'created': str(core_props.created) if hasattr(core_props, 'created') else "",
                    'modified': str(core_props.modified) if hasattr(core_props, 'modified') else "",
                    'last_modified_by': core_props.last_modified_by or ""
                })
            except Exception as prop_err:
                print(f"Warning: Could not extract all properties: {prop_err}")
            
            content.set_metadata(metadata)
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"--- Slide {slide_num + 1} ---\n"
                
                # Extract slide title if available
                if slide.shapes.title:
                    slide_text += f"Title: {slide.shapes.title.text}\n\n"
                
                # Process each shape in the slide
                for shape_idx, shape in enumerate(slide.shapes):
                    # Extract text from shape
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_text += para_text + "\n"
                    
                    # Extract tables
                    if hasattr(shape, 'table'):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text_frame:
                                    cell_content = cell.text_frame.text.strip()
                                    row_text.append(cell_content)
                            table_text.append(" | ".join(row_text))
                        
                        table_content = "\n".join(table_text)
                        content.add_table(table_content, slide_num, shape_idx)
                        slide_text += f"[Table content on this slide]\n"
                
                # Add the slide text
                content.add_text(slide_text, slide_num)
                
                # Process images in the slide
                img_index = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # Extract image from shape
                            image = shape.image
                            image_data = image.blob
                            
                            # Process the image
                            description = self.process_image(image_data)
                            content.add_image(image_data, description, slide_num, img_index)
                            img_index += 1
                        except Exception as img_err:
                            print(f"Error processing image in slide {slide_num + 1}: {img_err}")
            
            return content
            
        except Exception as e:
            print(f"Error reading PPTX file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content 