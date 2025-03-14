import os
import yaml
import json
from typing import Dict, Any, Tuple, Generator, List, Optional
from abc import ABC, abstractmethod
from langchain_ollama import ChatOllama
from langchain_groq import  ChatGroq
from langchain.chat_models import init_chat_model
# import docx
import PyPDF2
import markdown
import pptx
from bs4 import BeautifulSoup
import csv
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io

# TODO: json, csv



class DocumentContent:
    """Container for extracted document content with text, images, tables, etc."""
    
    def __init__(self):
        self.text_blocks = []
        self.images = []
        self.tables = []
        self.metadata = {}
    
    def add_text(self, text: str, page_num: Optional[int] = None):
        """Add a text block with optional page number."""
        if page_num is not None:
            self.text_blocks.append(f"--- Page {page_num + 1} ---\n{text}")
        else:
            self.text_blocks.append(text)
    
    def add_image(self, image_data: bytes, description: str, page_num: Optional[int] = None, img_index: int = 0):
        """Add an image with its description and location info."""
        self.images.append({
            "data": image_data,
            "description": description,
            "page": page_num,
            "index": img_index
        })
        
        if page_num is not None:
            self.text_blocks.append(f"--- Image on page {page_num + 1}, #{img_index + 1} ---\n{description}")
        else:
            self.text_blocks.append(f"--- Image #{img_index + 1} ---\n{description}")
    
    def add_table(self, table_text: str, page_num: Optional[int] = None, table_index: int = 0):
        """Add a table with its content and location info."""
        self.tables.append({
            "content": table_text,
            "page": page_num,
            "index": table_index
        })
        
        if page_num is not None:
            self.text_blocks.append(f"--- Table on page {page_num + 1}, #{table_index + 1} ---\n{table_text}")
        else:
            self.text_blocks.append(f"--- Table #{table_index + 1} ---\n{table_text}")
    
    def get_text(self) -> str:
        """Get all content as a single text string."""
        return "\n\n".join(self.text_blocks)
    
    def set_metadata(self, metadata: Dict[str, Any]):
        """Set document metadata."""
        self.metadata = metadata
    
    def update_metadata(self, metadata: Dict[str, Any]):
        """Update document metadata with new values."""
        self.metadata.update(metadata)
    


class DocumentReader(ABC):
    """Base class for document type readers."""
    
    def __init__(self, llm, vision_llm=None):
        self.llm = llm
        self.vision_llm = vision_llm
    
    @abstractmethod
    def read(self, file_path: str) -> DocumentContent:
        """Read document and extract content."""
        pass
    
    def process_image(self, image_data: bytes) -> str:
        """Process image data and return description."""
        if not image_data:
            return "[Image: No data available]"
        
        try:
            
            if not self.vision_llm:
                # Use OCR as fallback if no vision LLM
                img = Image.open(io.BytesIO(image_data))
                text = pytesseract.image_to_string(img)
                return f"[Image content: {text or 'No text detected'}]"
            
            # Use vision LLM to describe the image
            description = self.vision_llm.invoke(f"Describe this image in detail", image_data)
            return f"[Image description: {description}]"
        except Exception as e:
            print(f"Error processing image: {e}")
            return "[Image: Unable to process]"



class TextReader(DocumentReader):
    """Reader for plain text files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
                content.add_text(text)
                
                # Basic metadata
                content.set_metadata({
                    'source': file_path,
                    'format': 'txt',
                    'size_bytes': os.path.getsize(file_path)
                })
                
            return content
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content
        

class PDFReader(DocumentReader):
    """Reader for PDF files using PyMuPDF."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        
        try:
            # Open the PDF with PyMuPDF
            doc = fitz.open(file_path)
            
            # Extract document metadata
            pdf_metadata = {
                "source": file_path,
                "format": "pdf",
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
                "keywords": doc.metadata.get("keywords", ""),
                "creator": doc.metadata.get("creator", ""),
                "producer": doc.metadata.get("producer", ""),
                "creationDate": doc.metadata.get("creationDate", ""),
                "modDate": doc.metadata.get("modDate", ""),
                "pageCount": len(doc),
                "size_bytes": os.path.getsize(file_path)
            }
            
            content.set_metadata(pdf_metadata)
            
            # Process each page
            for page_num, page in enumerate(doc):
                # Extract text
                page_text = page.get_text("text")
                content.add_text(page_text, page_num)
                
                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]  # get the XREF of the image
                    base_image = doc.extract_image(xref)
                    image_data = base_image["image"]
                    
                    # Process the image
                    img_description = self.process_image(image_data)
                    content.add_image(image_data, img_description, page_num, img_index)
                
                # Try to detect tables (simplified approach)
                blocks = page.get_text("dict")["blocks"]
                table_index = 0
                for b in blocks:
                    if "lines" in b and len(b["lines"]) > 2:
                        # Check if this might be a table based on structure
                        rect = fitz.Rect(b["bbox"])
                        lines_count = len(b["lines"])
                        spans_per_line = [len(line["spans"]) for line in b["lines"]]
                        
                        # Simple heuristic: if multiple lines with multiple spans each
                        if lines_count > 2 and max(spans_per_line) > 2:
                            # Extract table text from the specified area
                            table_text = page.get_text("text", clip=rect)
                            content.add_table(f"[Table content:\n{table_text}\n]", page_num, table_index)
                            table_index += 1
            
            # Close the document
            doc.close()
            
            return content
        
        except Exception as e:
            print(f"Error reading PDF file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content

class MarkdownReader(DocumentReader):
    """Reader for Markdown files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                md_content = f.read()
                html = markdown.markdown(md_content)
                soup = BeautifulSoup(html, 'html.parser')
                text = soup.get_text()
                content.add_text(text)
                
                # Extract images if any
                img_index = 0
                for img in soup.find_all('img'):
                    src = img.get('src', '')
                    alt = img.get('alt', 'No description')
                    
                    # For local images referenced in the markdown
                    if src and not src.startswith(('http://', 'https://')):
                        img_path = os.path.join(os.path.dirname(file_path), src)
                        if os.path.exists(img_path):
                            try:
                                with open(img_path, 'rb') as img_file:
                                    img_data = img_file.read()
                                
                                description = self.process_image(img_data) if self.vision_llm else f"[Image: {alt}]"
                                content.add_image(img_data, description, None, img_index)
                                img_index += 1
                            except Exception as img_err:
                                print(f"Error processing image {src}: {img_err}")
                    else:
                        content.add_text(f"[External image: {src}, Description: {alt}]")
                
                # Basic metadata
                content.set_metadata({
                    'source': file_path,
                    'format': 'md',
                    'size_bytes': os.path.getsize(file_path)
                })
                
            return content
        except Exception as e:
            print(f"Error reading markdown file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content
        
class HTMLReader(DocumentReader):
    """Reader for HTML files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html_content = f.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                text = soup.get_text()
                content.add_text(text)
                
                # Extract images if any
                img_index = 0
                for img in soup.find_all('img'):
                    src = img.get('src', '')
                    alt = img.get('alt', 'No description')
                    
                    # For local images referenced in the HTML
                    if src and not src.startswith(('http://', 'https://')):
                        img_path = os.path.join(os.path.dirname(file_path), src)
                        if os.path.exists(img_path):
                            try:
                                with open(img_path, 'rb') as img_file:
                                    img_data = img_file.read()
                                
                                description = self.process_image(img_data) if self.vision_llm else f"[Image: {alt}]"
                                content.add_image(img_data, description, None, img_index)
                                img_index += 1
                            except Exception as img_err:
                                print(f"Error processing image {src}: {img_err}")
                    else:
                        content.add_text(f"[External image: {src}, Description: {alt}]")
                
                # Extract tables
                table_index = 0
                for table in soup.find_all('table'):
                    table_text = table.get_text()
                    content.add_table(f"[Table content:\n{table_text}\n]", None, table_index)
                    table_index += 1
                
                # Extract metadata from head tags
                meta_data = {
                    'source': file_path,
                    'format': 'html',
                    'size_bytes': os.path.getsize(file_path),
                    'title': soup.title.string if soup.title else ""
                }
                
                # Extract other meta tags
                for meta in soup.find_all('meta'):
                    name = meta.get('name', meta.get('property', ''))
                    if name and meta.get('content'):
                        meta_data[name] = meta.get('content')
                
                content.set_metadata(meta_data)
                
            return content
        except Exception as e:
            print(f"Error reading HTML file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content



class JSONReader(DocumentReader):
    """Reader for JSON files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                json_data = json.load(f)
                
                # Format JSON prettily and add as text
                formatted_json = json.dumps(json_data, indent=2)
                content.add_text(formatted_json)
                
                # Try to extract any nested tables or structured data
                self._process_json_object(json_data, content)
                
                # Basic metadata
                content.set_metadata({
                    'source': file_path,
                    'format': 'json',
                    'size_bytes': os.path.getsize(file_path)
                })
                
            return content
        except Exception as e:
            print(f"Error reading JSON file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content
    
    def _process_json_object(self, data, content, table_index=0, path=""):
        """Process JSON objects to extract table-like structures."""
        if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
            # This looks like a table: list of dictionaries
            # First, get all possible keys from all objects
            all_keys = set()
            for item in data:
                all_keys.update(item.keys())
            
            # Create a table representation
            table_rows = [",".join(all_keys)]  # Header row
            
            for item in data:
                row = []
                for key in all_keys:
                    value = item.get(key, "")
                    # Handle nested structures by converting to string
                    if isinstance(value, (dict, list)):
                        value = json.dumps(value)
                    row.append(str(value))
                table_rows.append(",".join(row))
            
            table_text = "\n".join(table_rows)
            path_info = f" at {path}" if path else ""
            content.add_table(f"[JSON Array as Table{path_info}:\n{table_text}\n]", None, table_index)
            return table_index + 1
        
        # Process nested structures
        elif isinstance(data, dict):
            for key, value in data.items():
                new_path = f"{path}.{key}" if path else key
                if isinstance(value, (dict, list)):
                    table_index = self._process_json_object(value, content, table_index, new_path)
        
        elif isinstance(data, list):
            for i, item in enumerate(data):
                new_path = f"{path}[{i}]"
                if isinstance(item, (dict, list)):
                    table_index = self._process_json_object(item, content, table_index, new_path)
        
        return table_index


class CSVReader(DocumentReader):
    """Reader for CSV files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            # Try to detect the dialect first
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                sample = f.read(1024)
                sniffer = csv.Sniffer()
                has_header = sniffer.has_header(sample)
                dialect = sniffer.sniff(sample)
            
            # Now read the CSV with the detected dialect
            with open(file_path, 'r', encoding='utf-8', errors='replace', newline='') as f:
                reader = csv.reader(f, dialect=dialect)
                rows = list(reader)
                
                # Extract header row if detected
                header_row = None
                data_rows = rows
                if has_header and rows:
                    header_row = rows[0]
                    data_rows = rows[1:]
                
                # Calculate basic statistics
                row_count = len(data_rows)
                column_count = len(header_row) if header_row else (len(rows[0]) if rows else 0)
                
                # Format the CSV data as text
                formatted_rows = []
                if header_row:
                    formatted_rows.append(" | ".join(header_row))
                    formatted_rows.append("-" * (sum(len(h) + 3 for h in header_row)))
                
                for row in data_rows:
                    formatted_rows.append(" | ".join(row))
                
                formatted_csv = "\n".join(formatted_rows)
                content.add_text(formatted_csv)
                
                # Add as table as well
                content.add_table(formatted_csv)
                
                # Set metadata
                content.set_metadata({
                    'source': file_path,
                    'format': 'csv',
                    'size_bytes': os.path.getsize(file_path),
                    'rows': row_count + (1 if has_header else 0),
                    'columns': column_count,
                    'has_header': has_header
                })
                
                # Try to determine data types in each column
                if rows:
                    column_types = self._analyze_column_types(data_rows, column_count)
                    content.update_metadata({'column_types': column_types})
                
            return content
        except Exception as e:
            print(f"Error reading CSV file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content
    
    def _analyze_column_types(self, data_rows, column_count):
        """Analyze and determine likely data types for each column."""
        column_types = []
        
        for col_idx in range(column_count):
            # Get all values in this column
            column_values = [row[col_idx] if col_idx < len(row) else "" for row in data_rows]
            
            # Check if all values can be converted to numbers
            numeric_count = 0
            date_count = 0
            empty_count = 0
            
            for value in column_values:
                if not value.strip():
                    empty_count += 1
                    continue
                    
                # Try to convert to float
                try:
                    float(value)
                    numeric_count += 1
                    continue
                except ValueError:
                    pass
                
                # Try basic date format detection
                # This is a simplified approach - in a real app, would use more sophisticated date detection
                if len(value) >= 8 and (('/' in value) or ('-' in value) or ('.' in value)):
                    date_count += 1
            
            # Determine column type based on majority
            non_empty_count = len(column_values) - empty_count
            if non_empty_count == 0:
                column_type = "empty"
            elif numeric_count / non_empty_count > 0.8:
                column_type = "numeric"
            elif date_count / non_empty_count > 0.8:
                column_type = "date"
            else:
                column_type = "text"
                
            column_types.append(column_type)
            
        return column_types
    

class DocxReader(DocumentReader):
    """Reader for Microsoft Word (.docx) files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            import docx
            
            # Open the document
            doc = docx.Document(file_path)
            
            # Extract metadata
            doc_properties = doc.core_properties
            metadata = {
                'source': file_path,
                'format': 'docx',
                'title': doc_properties.title or "",
                'author': doc_properties.author or "",
                'subject': doc_properties.subject or "",
                'keywords': doc_properties.keywords or "",
                'created': str(doc_properties.created) if doc_properties.created else "",
                'modified': str(doc_properties.modified) if doc_properties.modified else "",
                'last_modified_by': doc_properties.last_modified_by or "",
                'size_bytes': os.path.getsize(file_path)
            }
            content.set_metadata(metadata)
            
            # Extract all paragraphs
            for para in doc.paragraphs:
                if para.text.strip():  # Only add non-empty paragraphs
                    content.add_text(para.text)
            
            # Extract tables
            table_index = 0
            for table in doc.tables:
                table_data = []
                for i, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        # Get text from each cell, handling potential nested content
                        cell_text = ""
                        for paragraph in cell.paragraphs:
                            if paragraph.text.strip():
                                cell_text += paragraph.text + " "
                        row_data.append(cell_text.strip())
                    table_data.append(" | ".join(row_data))
                
                table_text = "\n".join(table_data)
                content.add_table(table_text, None, table_index)
                table_index += 1
            
            # Handle images
            img_index = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                        image_data = image_part.blob
                        
                        # Process the image
                        description = self.process_image(image_data)
                        content.add_image(image_data, description, None, img_index)
                        img_index += 1
                    except Exception as img_err:
                        print(f"Error processing image in DOCX: {img_err}")
            
            return content
            
        except Exception as e:
            print(f"Error reading DOCX file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content


class PptxReader(DocumentReader):
    """Reader for Microsoft PowerPoint (.pptx) files."""
    
    def read(self, file_path: str) -> DocumentContent:
        content = DocumentContent()
        try:
            from pptx import Presentation
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract basic metadata
            metadata = {
                'source': file_path,
                'format': 'pptx',
                'slide_count': len(prs.slides),
                'size_bytes': os.path.getsize(file_path)
            }
            
            # Try to extract more detailed properties if available
            try:
                core_props = prs.core_properties
                metadata.update({
                    'title': core_props.title or "",
                    'author': core_props.author or "",
                    'subject': core_props.subject or "",
                    'keywords': core_props.keywords or "",
                    'created': str(core_props.created) if hasattr(core_props, 'created') else "",
                    'modified': str(core_props.modified) if hasattr(core_props, 'modified') else "",
                    'last_modified_by': core_props.last_modified_by or ""
                })
            except Exception as prop_err:
                print(f"Warning: Could not extract all properties: {prop_err}")
            
            content.set_metadata(metadata)
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"--- Slide {slide_num + 1} ---\n"
                
                # Extract slide title if available
                if slide.shapes.title:
                    slide_text += f"Title: {slide.shapes.title.text}\n\n"
                
                # Process each shape in the slide
                for shape_idx, shape in enumerate(slide.shapes):
                    # Extract text from shape
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_text += para_text + "\n"
                    
                    # Extract tables
                    if hasattr(shape, 'table'):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text_frame:
                                    cell_content = cell.text_frame.text.strip()
                                    row_text.append(cell_content)
                            table_text.append(" | ".join(row_text))
                        
                        table_content = "\n".join(table_text)
                        content.add_table(table_content, slide_num, shape_idx)
                        slide_text += f"[Table content on this slide]\n"
                
                # Add the slide text
                content.add_text(slide_text, slide_num)
                
                # Process images in the slide
                img_index = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # Extract image from shape
                            image = shape.image
                            image_data = image.blob
                            
                            # Process the image
                            description = self.process_image(image_data)
                            content.add_image(image_data, description, slide_num, img_index)
                            img_index += 1
                        except Exception as img_err:
                            print(f"Error processing image in slide {slide_num + 1}: {img_err}")
            
            return content
            
        except Exception as e:
            print(f"Error reading PPTX file {file_path}: {e}")
            content.add_text(f"Error reading file: {str(e)}")
            return content